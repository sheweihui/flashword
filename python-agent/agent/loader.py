"""文件加载器 — 可插拔架构，支持 PDF/PPTX/TXT/MD/DOCX。"""

from __future__ import annotations

from abc import ABC, abstractmethod
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Type


class Document:
    """加载器产出的统一文档对象。"""
    def __init__(self, text: str, metadata: Optional[Dict] = None):
        self.text = text
        self.metadata = metadata or {}

    def __repr__(self) -> str:
        return f"Document(text={len(self.text)} chars, metadata={self.metadata})"


# ============================================================
# 抽象基类
# ============================================================

class BaseLoader(ABC):
    """文档加载器抽象基类。

    所有格式加载器继承此类，实现 load() 方法，
    返回统一的 Document 对象。
    """

    @abstractmethod
    def load(self, file_path: str | Path) -> Document:
        """加载并解析文档。

        Args:
            file_path: 文件路径。

        Returns:
            Document 对象，包含 text 和 metadata。
        """
        ...

    @property
    @abstractmethod
    def supported_extensions(self) -> List[str]:
        """支持的文件扩展名列表（含.）。"""
        ...

    @staticmethod
    def _validate(file_path: str | Path) -> Path:
        path = Path(file_path).resolve()
        if not path.exists():
            raise FileNotFoundError(f"文件不存在: {path}")
        if not path.is_file():
            raise ValueError(f"不是文件: {path}")
        return path


# ============================================================
# PDF 加载器
# ============================================================

class PDFLoader(BaseLoader):
    """PDF 文件加载器。"""

    @property
    def supported_extensions(self) -> List[str]:
        return [".pdf"]

    def load(self, file_path: str | Path) -> Document:
        path = self._validate(file_path)
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text and text.strip():
                pages.append(f"--- 第 {i + 1} 页 ---\n{text.strip()}")

        return Document(
            text="\n\n".join(pages),
            metadata={
                "source_path": str(path),
                "file_name": path.name,
                "page_count": len(reader.pages),
                "format": "pdf",
            },
        )


# ============================================================
# PPTX 加载器
# ============================================================

class PPTXLoader(BaseLoader):
    """PPTX 文件加载器。"""

    @property
    def supported_extensions(self) -> List[str]:
        return [".pptx"]

    def load(self, file_path: str | Path) -> Document:
        path = self._validate(file_path)
        from pptx import Presentation

        prs = Presentation(str(path))
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            if texts:
                slides.append(f"--- 幻灯片 {i + 1} ---\n" + "\n".join(texts))

        return Document(
            text="\n\n".join(slides),
            metadata={
                "source_path": str(path),
                "file_name": path.name,
                "slide_count": len(prs.slides),
                "format": "pptx",
            },
        )


# ============================================================
# DOCX 加载器
# ============================================================

class DOCXLoader(BaseLoader):
    """DOCX 文件加载器。"""

    @property
    def supported_extensions(self) -> List[str]:
        return [".docx"]

    def load(self, file_path: str | Path) -> Document:
        path = self._validate(file_path)
        from docx import Document as DocxDocument

        doc = DocxDocument(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

        return Document(
            text="\n\n".join(paragraphs),
            metadata={
                "source_path": str(path),
                "file_name": path.name,
                "paragraph_count": len(paragraphs),
                "format": "docx",
            },
        )


# ============================================================
# TXT/MD 加载器
# ============================================================

class TextLoader(BaseLoader):
    """纯文本 / Markdown 文件加载器。"""

    @property
    def supported_extensions(self) -> List[str]:
        return [".txt", ".md", ".text"]

    def load(self, file_path: str | Path) -> Document:
        path = self._validate(file_path)
        text = path.read_text(encoding="utf-8", errors="replace")

        return Document(
            text=text.strip(),
            metadata={
                "source_path": str(path),
                "file_name": path.name,
                "file_size": path.stat().st_size,
                "format": path.suffix[1:] if path.suffix else "txt",
            },
        )


# ============================================================
# Loader 工厂
# ============================================================

class LoaderFactory:
    """加载器工厂 — 根据文件扩展名自动选择加载器。"""

    def __init__(self):
        self._loaders: Dict[str, Type[BaseLoader]] = {}

    def register(self, loader_class: Type[BaseLoader]) -> None:
        """注册一个加载器。"""
        loader = loader_class()
        for ext in loader.supported_extensions:
            self._loaders[ext.lower()] = loader_class

    def create(self, file_path: str | Path) -> BaseLoader:
        """根据文件扩展名创建对应的加载器。

        Args:
            file_path: 文件路径。

        Returns:
            BaseLoader 实例。

        Raises:
            ValueError: 如果不支持该格式。
        """
        suffix = Path(file_path).suffix.lower()
        loader_class = self._loaders.get(suffix)
        if loader_class is None:
            supported = ", ".join(sorted(self._loaders.keys()))
            raise ValueError(
                f"不支持的文件格式 '{suffix}'。支持: {supported}"
            )
        return loader_class()

    def list_supported(self) -> List[str]:
        """列出所有支持的扩展名。"""
        return sorted(self._loaders.keys())

    def load(self, file_path: str | Path) -> Document:
        """一键加载：自动识别格式并解析。"""
        loader = self.create(file_path)
        return loader.load(file_path)


# ============================================================
# 便捷函数（保持兼容）
# ============================================================

def get_title_from_filename(file_path: str | Path) -> str:
    return Path(file_path).stem


def _create_default_factory() -> LoaderFactory:
    """创建默认工厂并注册所有内置加载器。"""
    factory = LoaderFactory()
    factory.register(PDFLoader)
    factory.register(PPTXLoader)
    factory.register(DOCXLoader)
    factory.register(TextLoader)
    return factory


_default_factory = None


def load_file(file_path: str | Path) -> Document:
    """便捷函数 — 自动识别格式并加载文件。"""
    global _default_factory
    if _default_factory is None:
        _default_factory = _create_default_factory()
    return _default_factory.load(file_path)


def list_supported_formats() -> List[str]:
    """列出所有支持的格式。"""
    global _default_factory
    if _default_factory is None:
        _default_factory = _create_default_factory()
    return _default_factory.list_supported()
